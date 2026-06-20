import fitz  # PyMuPDF
import io
import base64
import csv
import json
import os
import uuid
from docx import Document

IMAGES_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.dirname(__file__))),
    "data", "images"
)
os.makedirs(IMAGES_DIR, exist_ok=True)

class DocumentParser:
    @staticmethod
    def parse_pdf_advanced(file_bytes: bytes):
        """Extracts text, tracks true page boundaries, and renders embedded images."""
        text = ""
        extracted_images = []
        
        # Open PDF from system memory
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        total_pages = len(doc)
        
        for page_num in range(total_pages):
            page = doc.load_page(page_num)
            text += f"\n--- [Page {page_num + 1}] ---\n" + page.get_text()
            
            # Extract image assets
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                # Save image to disk
                image_id = str(uuid.uuid4())
                filename = f"{image_id}.{image_ext}"
                filepath = os.path.join(IMAGES_DIR, filename)
                with open(filepath, "wb") as f:
                    f.write(image_bytes)
                
                # Inject marker directly into the text for this page
                img_url = f"/api/images/{filename}"
                text += f"\n\n[Image available at: {img_url}]\n\n"
                extracted_images.append(img_url)
                
        return {
            "text": text.strip(),
            "page_count": total_pages,
            "images": extracted_images
        }

    @staticmethod
    def parse_docx_advanced(file_bytes: bytes):
        """Extracts text contents and logs default structural layouts from DOCX."""
        text = []
        docx_file = io.BytesIO(file_bytes)
        doc = Document(docx_file)
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text)
        
        # Also extract table content
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text.append(row_text)
                
        full_text = "\n".join(text)
        # Estimate pages for docx files based on average word counts (approx 400 words per page)
        estimated_pages = max(1, len(full_text.split()) // 400)
        
        return {
            "text": full_text.strip(),
            "page_count": estimated_pages,
            "images": []  # Base docx image matching requires separate zipfile parsing
        }

    @staticmethod
    def parse_excel(file_bytes: bytes):
        """Extracts text from all sheets in an Excel (.xlsx/.xls) file."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            sections = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_vals = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_vals):
                        rows.append(" | ".join(row_vals))
                if rows:
                    sections.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
            full_text = "\n\n".join(sections)
            page_count = max(1, len(wb.sheetnames))
            return {
                "text": full_text.strip(),
                "page_count": page_count,
                "images": []
            }
        except Exception as exc:
            raise ValueError(f"Failed to parse Excel file: {exc}")

    @staticmethod
    def parse_csv(file_bytes: bytes):
        """Extracts all rows from a CSV file as readable text."""
        try:
            content = file_bytes.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(content))
            rows = []
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows.append(" | ".join(row))
            full_text = "\n".join(rows)
            estimated_pages = max(1, len(rows) // 50)
            return {
                "text": full_text.strip(),
                "page_count": estimated_pages,
                "images": []
            }
        except Exception as exc:
            raise ValueError(f"Failed to parse CSV file: {exc}")

    @staticmethod
    def parse_txt(file_bytes: bytes):
        """Reads plain text or markdown files directly."""
        try:
            content = file_bytes.decode("utf-8", errors="replace")
            words = content.split()
            estimated_pages = max(1, len(words) // 400)
            return {
                "text": content.strip(),
                "page_count": estimated_pages,
                "images": []
            }
        except Exception as exc:
            raise ValueError(f"Failed to parse text file: {exc}")

    @staticmethod
    def parse_pptx(file_bytes: bytes):
        """Extracts text from each slide of a PowerPoint presentation."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                if slide_parts:
                    slides_text.append(f"--- [Slide {slide_num}] ---\n" + "\n".join(slide_parts))
            full_text = "\n\n".join(slides_text)
            return {
                "text": full_text.strip(),
                "page_count": len(prs.slides),
                "images": []
            }
        except Exception as exc:
            raise ValueError(f"Failed to parse PowerPoint file: {exc}")

    @staticmethod
    def parse_ipynb(file_bytes: bytes):
        """Extracts code cells and markdown cells from a Jupyter Notebook."""
        try:
            nb = json.loads(file_bytes.decode("utf-8", errors="replace"))
            cells_text = []
            for cell in nb.get("cells", []):
                cell_type = cell.get("cell_type", "")
                source = "".join(cell.get("source", []))
                if not source.strip():
                    continue
                if cell_type == "markdown":
                    cells_text.append(f"[Markdown]\n{source}")
                elif cell_type == "code":
                    cells_text.append(f"[Code]\n{source}")
                    # Also include output text if available
                    for output in cell.get("outputs", []):
                        if output.get("output_type") in ("stream", "display_data", "execute_result"):
                            out_text = "".join(output.get("text", []))
                            if out_text.strip():
                                cells_text.append(f"[Output]\n{out_text}")
            full_text = "\n\n".join(cells_text)
            estimated_pages = max(1, len(full_text.split()) // 400)
            return {
                "text": full_text.strip(),
                "page_count": estimated_pages,
                "images": []
            }
        except Exception as exc:
            raise ValueError(f"Failed to parse Jupyter Notebook: {exc}")